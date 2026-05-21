#!/usr/bin/env python3
"""Extract teaching materials into a chunked JSONL corpus."""

from __future__ import annotations

import argparse
import json
import re
from dataclasses import dataclass, field
from pathlib import Path

SUPPORTED_EXTENSIONS = {".txt", ".md", ".pdf", ".docx", ".pptx"}
TYPE_BY_EXTENSION = {
    ".txt": "txt",
    ".md": "md",
    ".pdf": "pdf",
    ".docx": "docx",
    ".pptx": "pptx",
}
MISSING_MESSAGES = {
    "pdf": "Skipping PDF extraction; install pymupdf.",
    "docx": "Skipping DOCX extraction; install python-docx.",
    "pptx": "Skipping PPTX extraction; install python-pptx.",
}


@dataclass
class ExtractStats:
    files_scanned: int = 0
    files_processed: int = 0
    files_skipped: int = 0
    unsupported: int = 0
    chunks_written: int = 0
    missing_dependencies: set[str] = field(default_factory=set)


def normalize_text(text: str) -> str:
    text = text.replace("\x00", " ")
    text = re.sub(r"[ \t\r\f\v]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def chunk_text(text: str, chunk_chars: int, overlap_chars: int) -> list[str]:
    text = normalize_text(text)
    if not text:
        return []
    if len(text) <= chunk_chars:
        return [text]

    chunks: list[str] = []
    start = 0
    while start < len(text):
        end = min(len(text), start + chunk_chars)
        if end < len(text):
            boundary = max(text.rfind("\n", start, end), text.rfind(". ", start, end))
            if boundary > start + int(chunk_chars * 0.55):
                end = boundary + 1
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end >= len(text):
            break
        start = max(0, end - overlap_chars)
    return chunks


def safe_stem(path: Path) -> str:
    return re.sub(r"[^A-Za-z0-9]+", "_", path.stem).strip("_").lower() or "source"


def is_ignored_path(path: Path) -> bool:
    return (
        any(part.startswith(".") for part in path.parts)
        or path.name == ".gitkeep"
        or path.name.startswith("~$")
    )


def read_text_with_fallback(path: Path) -> str:
    for encoding in ["utf-8", "utf-8-sig", "cp1250", "latin-1"]:
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    return path.read_text(encoding="utf-8", errors="replace")


def read_txt_like(path: Path) -> list[dict[str, object]]:
    return [{"page": None, "slide": None, "text": read_text_with_fallback(path)}]


def read_pdf(path: Path, stats: ExtractStats) -> list[dict[str, object]]:
    try:
        import fitz  # type: ignore
    except ImportError:
        stats.missing_dependencies.add("pdf")
        return []

    pages: list[dict[str, object]] = []
    with fitz.open(path) as document:
        for index, page in enumerate(document, start=1):
            pages.append({"page": index, "slide": None, "text": page.get_text("text")})
    return pages


def read_docx(path: Path, stats: ExtractStats) -> list[dict[str, object]]:
    try:
        import docx  # type: ignore
    except ImportError:
        stats.missing_dependencies.add("docx")
        return []

    document = docx.Document(path)
    text = "\n".join(paragraph.text for paragraph in document.paragraphs)
    return [{"page": None, "slide": None, "text": text}]


def read_pptx(path: Path, stats: ExtractStats) -> list[dict[str, object]]:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        stats.missing_dependencies.add("pptx")
        return []

    presentation = Presentation(path)
    slides: list[dict[str, object]] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                parts.append(text)
        slides.append({"page": None, "slide": index, "text": "\n".join(parts)})
    return slides


def extract_source(path: Path, source_type: str, stats: ExtractStats) -> list[dict[str, object]]:
    if source_type in {"txt", "md"}:
        return read_txt_like(path)
    if source_type == "pdf":
        return read_pdf(path, stats)
    if source_type == "docx":
        return read_docx(path, stats)
    if source_type == "pptx":
        return read_pptx(path, stats)
    return []


def iter_material_files(materials_dir: Path, stats: ExtractStats) -> list[Path]:
    paths: list[Path] = []
    if not materials_dir.exists():
        return paths
    for path in sorted(materials_dir.rglob("*")):
        if not path.is_file() or is_ignored_path(path):
            continue
        stats.files_scanned += 1
        if path.suffix.lower() not in SUPPORTED_EXTENSIONS:
            stats.files_skipped += 1
            stats.unsupported += 1
            continue
        paths.append(path)
    return paths


def chunk_marker(page: object, slide: object) -> str:
    if isinstance(slide, int):
        return f"s{slide:03d}"
    if isinstance(page, int):
        return f"p{page:03d}"
    return "p000"


def build_corpus(materials_dir: Path, chunk_chars: int, overlap_chars: int) -> tuple[list[dict[str, object]], ExtractStats]:
    stats = ExtractStats()
    rows: list[dict[str, object]] = []
    for source_path in iter_material_files(materials_dir, stats):
        source_type = TYPE_BY_EXTENSION[source_path.suffix.lower()]
        extracted = extract_source(source_path, source_type, stats)
        if not extracted:
            stats.files_skipped += 1
            continue

        wrote_for_file = False
        for unit in extracted:
            chunks = chunk_text(str(unit.get("text", "")), chunk_chars, overlap_chars)
            for index, chunk in enumerate(chunks, start=1):
                page = unit.get("page")
                slide = unit.get("slide")
                marker = chunk_marker(page, slide)
                chunk_id = f"{safe_stem(source_path)}_{marker}_c{index:03d}"
                rows.append(
                    {
                        "chunk_id": chunk_id,
                        "source": source_path.name,
                        "source_path": source_path.as_posix(),
                        "source_type": source_type,
                        "page": page,
                        "slide": slide,
                        "text": chunk,
                        "char_count": len(chunk),
                    }
                )
                wrote_for_file = True
        if wrote_for_file:
            stats.files_processed += 1
        else:
            stats.files_skipped += 1
    stats.chunks_written = len(rows)
    return rows, stats


def main() -> int:
    parser = argparse.ArgumentParser(description="Extract teaching materials into data/generated/corpus.jsonl.")
    parser.add_argument("--materials-dir", type=Path, default=Path("materials"))
    parser.add_argument("--output", type=Path, default=Path("data/generated/corpus.jsonl"))
    parser.add_argument("--chunk-chars", type=int, default=3500)
    parser.add_argument("--overlap-chars", type=int, default=300)
    args = parser.parse_args()

    if args.chunk_chars < 500:
        parser.error("--chunk-chars must be at least 500")
    if args.overlap_chars < 0 or args.overlap_chars >= args.chunk_chars:
        parser.error("--overlap-chars must be non-negative and smaller than --chunk-chars")

    rows, stats = build_corpus(args.materials_dir, args.chunk_chars, args.overlap_chars)
    args.output.parent.mkdir(parents=True, exist_ok=True)
    with args.output.open("w", encoding="utf-8") as handle:
        for row in rows:
            handle.write(json.dumps(row, ensure_ascii=False) + "\n")

    for key in sorted(stats.missing_dependencies):
        print(MISSING_MESSAGES[key])
    print(f"files scanned: {stats.files_scanned}")
    print(f"files processed: {stats.files_processed}")
    print(f"files skipped: {stats.files_skipped}")
    print(f"unsupported files skipped: {stats.unsupported}")
    print(
        "missing optional dependencies: "
        + (", ".join(sorted(stats.missing_dependencies)) if stats.missing_dependencies else "none")
    )
    print(f"chunks written: {stats.chunks_written}")
    print(f"output: {args.output}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
